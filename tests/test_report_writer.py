"""
tests/test_report_writer.py — C2: Full Report + PowerPoint Generator

All tests run without a real API key.
Mock pattern: patch("ai.report_writer.call_llm")
"""

from __future__ import annotations

import json
import sys
from dataclasses import dataclass, field
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

VALID_CONCLUSION_JSON = json.dumps({
    "section_text": "This study identified 14 clusters in liver HCC tissue. "
                    "CD8+ T cells showed PDCD1 upregulation consistent with exhaustion.",
    "key_findings": [
        "14 clusters identified at resolution 0.5",
        "PDCD1 upregulated in CD8+ T cells (log2FC=2.1)",
        "Macrophage subpopulation enriched for immunosuppressive pathway",
    ],
    "cited_evidence": ["resolution=0.5", "PDCD1", "log2FC=2.1", "n_cells=5000"],
    "reasoning": "Assembled from analysis summary and narrative blocks.",
})

VALID_PERSPECTIVES_JSON = json.dumps({
    "section_text": "Future work should validate the exhausted CD8+ T cell population "
                    "by CyTOF and functional assays. The immunosuppressive macrophage "
                    "subpopulation represents a potential therapeutic target.",
    "key_findings": [],
    "cited_evidence": ["PDCD1", "macrophage"],
    "reasoning": "Derived from coherence review flags and narrative interpretation block.",
})

ANALYSIS_SUMMARY = {
    "study_context": {
        "tissue": "liver",
        "disease": "HCC",
        "n_cells": 5000,
        "n_batches": 3,
        "modalities": "RNA",
    },
    "qc_decisions": {
        "min_genes": 300,
        "max_mt_pct": 20,
        "cells_removed_pct": 12.3,
    },
    "clustering": {
        "resolution": 0.5,
        "n_clusters": 14,
        "silhouette_score": 0.41,
    },
    "cell_types": [
        {"cluster": 0, "cell_type": "CD8+ T cell", "confidence": "high", "n_cells": 412},
        {"cluster": 1, "cell_type": "Macrophage",  "confidence": "high", "n_cells": 380},
    ],
    "top_degs": [
        {"comparison": "tumour_vs_normal", "cluster": 0, "gene": "PDCD1", "log2fc": 2.1},
        {"comparison": "tumour_vs_normal", "cluster": 1, "gene": "CD68",  "log2fc": 1.8},
    ],
    "pathways": [
        {"cluster": 0, "pathway": "T cell exhaustion",      "padj": 0.001},
        {"cluster": 1, "pathway": "Innate immune response", "padj": 0.003},
    ],
}

STUDY_CONTEXT = {
    "dataset": {"name": "GSE166635", "tissue": "liver"},
    "disease": {"context": "hepatocellular carcinoma"},
    "biological_question": "Characterise the TME of HCC.",
    "objectives": ["Identify major cell types", "Find DEGs"],
}


def _config(ai_on: bool = True, module_on: bool = True) -> dict:
    modules = {"report_writer": module_on}
    return {
        "ai": {
            "features": ai_on,
            "provider": "ollama",
            "model": "llama3",
            "modules": modules,
        }
    }


def _mock_adata(summary: dict | None = None) -> MagicMock:
    adata = MagicMock()
    adata.uns = {"analysis_summary": summary or ANALYSIS_SUMMARY}
    return adata


# ---------------------------------------------------------------------------
# Fake NarrativeResult for speaker-notes tests
# ---------------------------------------------------------------------------

@dataclass
class _FakeBlock:
    block_name: str = ""
    narrative_text: str = ""
    cited_evidence: list[str] = field(default_factory=list)
    groundedness_score: float = 0.0


@dataclass
class _FakeNarrativeResult:
    blocks: list[_FakeBlock] = field(default_factory=list)
    timestamp: str = "2026-01-01T00:00:00+00:00"
    model: str = "llama3"
    provider: str = "ollama"
    skill_name: str = "narrative_generator"
    skill_version: str = "1.0"
    reasoning: str = ""
    overall_groundedness: float = 0.9


def _fake_narrative() -> _FakeNarrativeResult:
    return _FakeNarrativeResult(
        blocks=[
            _FakeBlock("qc_rationale",          "QC removed 12.3% of cells."),
            _FakeBlock("cell_type_landscape",    "14 clusters found including CD8+ T cells."),
            _FakeBlock("differential_expression","PDCD1 upregulated in exhausted T cells."),
            _FakeBlock("interpretation",         "Findings suggest immunosuppressive microenvironment."),
        ]
    )


# ---------------------------------------------------------------------------
# Helper — side_effect that returns conclusion then perspectives JSON
# ---------------------------------------------------------------------------

def _two_call_side_effect(*args, **kwargs):
    """First call → conclusion JSON, second call → perspectives JSON."""
    if not hasattr(_two_call_side_effect, "_count"):
        _two_call_side_effect._count = 0
    _two_call_side_effect._count += 1
    if _two_call_side_effect._count % 2 == 1:
        return VALID_CONCLUSION_JSON
    return VALID_PERSPECTIVES_JSON


def _reset_counter():
    if hasattr(_two_call_side_effect, "_count"):
        del _two_call_side_effect._count


# ---------------------------------------------------------------------------
# Tests — config gate
# ---------------------------------------------------------------------------

class TestConfigGate:
    def test_returns_none_when_ai_features_false(self, tmp_path):
        result = _run_with_mock(
            config=_config(ai_on=False),
            tmp_path=tmp_path,
        )
        assert result is None

    def test_returns_none_when_module_off(self, tmp_path):
        result = _run_with_mock(
            config=_config(module_on=False),
            tmp_path=tmp_path,
        )
        assert result is None

    def test_returns_none_when_runtime_ai_false(self, tmp_path):
        with patch("ai.report_writer.call_llm", return_value=VALID_CONCLUSION_JSON):
            from ai.report_writer import run
            result = run(
                _mock_adata(),
                _config(),
                STUDY_CONTEXT,
                report_dir=str(tmp_path),
                runtime_ai=False,
            )
        assert result is None


# ---------------------------------------------------------------------------
# Tests — happy path
# ---------------------------------------------------------------------------

class TestHappyPath:
    def setup_method(self):
        _reset_counter()

    def test_returns_report_writer_result(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        from ai.report_writer import ReportWriterResult
        assert isinstance(result, ReportWriterResult)

    def test_two_sections_returned(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        assert len(result.sections) == 2

    def test_section_names(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        names = [s.section_name for s in result.sections]
        assert "conclusion" in names
        assert "perspectives" in names

    def test_conclusion_has_key_findings(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        conclusion = next(s for s in result.sections if s.section_name == "conclusion")
        assert len(conclusion.key_findings) == 3

    def test_conclusion_has_section_text(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        conclusion = next(s for s in result.sections if s.section_name == "conclusion")
        assert "PDCD1" in conclusion.section_text

    def test_perspectives_has_section_text(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        perspectives = next(s for s in result.sections if s.section_name == "perspectives")
        assert len(perspectives.section_text) > 10

    def test_two_llm_calls_made(self, tmp_path):
        with patch("ai.report_writer.call_llm", side_effect=_two_call_side_effect) as mock_llm:
            from ai.report_writer import run
            run(_mock_adata(), _config(), STUDY_CONTEXT, report_dir=str(tmp_path))
        assert mock_llm.call_count == 2

    def test_ai_result_base_fields_populated(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        assert result.model == "llama3"
        assert result.provider == "ollama"
        assert result.skill_name == "report_writer"
        assert result.skill_version == "1.0"
        assert result.timestamp  # non-empty string


# ---------------------------------------------------------------------------
# Tests — file output
# ---------------------------------------------------------------------------

class TestFileOutput:
    def setup_method(self):
        _reset_counter()

    def test_conclusion_md_written(self, tmp_path):
        _run_with_mock(tmp_path=tmp_path)
        assert (tmp_path / "conclusion.md").exists()

    def test_perspectives_md_written(self, tmp_path):
        _run_with_mock(tmp_path=tmp_path)
        assert (tmp_path / "perspectives.md").exists()

    def test_conclusion_md_contains_key_findings_as_bullets(self, tmp_path):
        _run_with_mock(tmp_path=tmp_path)
        content = (tmp_path / "conclusion.md").read_text()
        assert "- " in content  # markdown bullets

    def test_conclusion_md_has_section_text(self, tmp_path):
        _run_with_mock(tmp_path=tmp_path)
        content = (tmp_path / "conclusion.md").read_text()
        assert "PDCD1" in content

    def test_report_dir_created_if_missing(self, tmp_path):
        subdir = tmp_path / "nested" / "output"
        assert not subdir.exists()
        _run_with_mock(tmp_path=subdir)
        assert subdir.exists()


# ---------------------------------------------------------------------------
# Tests — PowerPoint
# ---------------------------------------------------------------------------

class TestPowerPoint:
    def setup_method(self):
        _reset_counter()

    def test_pptx_written_to_report_dir(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        assert (tmp_path / "presentation.pptx").exists()

    def test_pptx_path_in_result(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path)
        assert result.pptx_path is not None
        assert result.pptx_path.endswith("presentation.pptx")

    def test_pptx_has_eight_slides(self, tmp_path):
        _run_with_mock(tmp_path=tmp_path)
        from pptx import Presentation
        prs = Presentation(str(tmp_path / "presentation.pptx"))
        assert len(prs.slides) == 8

    def test_all_slides_have_speaker_notes(self, tmp_path):
        _run_with_mock(tmp_path=tmp_path)
        from pptx import Presentation
        prs = Presentation(str(tmp_path / "presentation.pptx"))
        for i, slide in enumerate(prs.slides):
            notes_text = slide.notes_slide.notes_text_frame.text
            assert notes_text.strip(), f"Slide {i+1} has empty speaker notes"

    def test_pptx_path_none_when_pptx_unavailable(self, tmp_path):
        """If python-pptx import fails, pptx_path is None but sections still written."""
        _reset_counter()
        with patch("ai.report_writer.call_llm", side_effect=_two_call_side_effect):
            with patch("ai.report_writer._build_pptx", side_effect=ImportError("python-pptx")):
                from ai.report_writer import run
                result = run(
                    _mock_adata(), _config(), STUDY_CONTEXT, report_dir=str(tmp_path)
                )
        assert result is not None
        assert result.pptx_path is None
        # Markdown sections still written
        assert (tmp_path / "conclusion.md").exists()

    def test_missing_figures_dir_handled(self, tmp_path):
        """figures_dir pointing to non-existent path does not crash."""
        result = _run_with_mock(
            tmp_path=tmp_path,
            figures_dir=str(tmp_path / "no_such_figures"),
        )
        assert (tmp_path / "presentation.pptx").exists()

    def test_pptx_with_narrative_result_uses_speaker_notes(self, tmp_path):
        """Speaker notes populated from narrative blocks when narrative_result provided."""
        _reset_counter()
        narrative = _fake_narrative()
        with patch("ai.report_writer.call_llm", side_effect=_two_call_side_effect):
            from ai.report_writer import run
            result = run(
                _mock_adata(), _config(), STUDY_CONTEXT,
                narrative_result=narrative,
                report_dir=str(tmp_path),
            )
        from pptx import Presentation
        prs = Presentation(str(tmp_path / "presentation.pptx"))
        # Slide 3 (UMAP) should have cell_type_landscape narrative text
        notes_slide3 = prs.slides[2].notes_slide.notes_text_frame.text
        assert "14 clusters" in notes_slide3


# ---------------------------------------------------------------------------
# Tests — graceful degradation
# ---------------------------------------------------------------------------

class TestGracefulDegradation:
    def setup_method(self):
        _reset_counter()

    def test_narrative_result_none_still_generates_sections(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path, narrative_result=None)
        assert len(result.sections) == 2

    def test_coherence_review_none_handled(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path, coherence_review=None)
        assert result is not None

    def test_cluster_annotations_none_handled(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path, cluster_annotations=None)
        assert result is not None

    def test_deg_validation_none_handled(self, tmp_path):
        result = _run_with_mock(tmp_path=tmp_path, deg_validation=None)
        assert result is not None

    def test_empty_analysis_summary_handled(self, tmp_path):
        """adata with no analysis_summary in uns generates report without crashing."""
        adata = _mock_adata(summary={})
        _reset_counter()
        with patch("ai.report_writer.call_llm", side_effect=_two_call_side_effect):
            from ai.report_writer import run
            result = run(adata, _config(), STUDY_CONTEXT, report_dir=str(tmp_path))
        assert result is not None

    def test_llm_returns_invalid_json_handled(self, tmp_path):
        """Bad JSON from LLM does not crash — raw text used as section_text."""
        with patch("ai.report_writer.call_llm", return_value="this is not json"):
            from ai.report_writer import run
            result = run(
                _mock_adata(), _config(), STUDY_CONTEXT, report_dir=str(tmp_path)
            )
        assert result is not None
        assert len(result.sections) == 2

    def test_adata_none_handled(self, tmp_path):
        """adata=None uses empty analysis_summary."""
        _reset_counter()
        with patch("ai.report_writer.call_llm", side_effect=_two_call_side_effect):
            from ai.report_writer import run
            result = run(None, _config(), STUDY_CONTEXT, report_dir=str(tmp_path))
        assert result is not None


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _run_with_mock(
    config=None,
    tmp_path=None,
    narrative_result=None,
    coherence_review=None,
    cluster_annotations=None,
    deg_validation=None,
    figures_dir=None,
):
    """Convenience wrapper: runs report_writer.run with mocked call_llm."""
    if config is None:
        config = _config()
    _reset_counter()
    with patch("ai.report_writer.call_llm", side_effect=_two_call_side_effect):
        from ai.report_writer import run
        return run(
            _mock_adata(),
            config,
            STUDY_CONTEXT,
            narrative_result=narrative_result,
            coherence_review=coherence_review,
            cluster_annotations=cluster_annotations,
            deg_validation=deg_validation,
            report_dir=str(tmp_path),
            figures_dir=figures_dir,
        )
