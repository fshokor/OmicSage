"""
ai/report_writer.py — C2: Full Report + PowerPoint Generator

Two deliverables:
  1. AI-generated report sections (conclusion.md, perspectives.md) via LLM
  2. Rule-based 8-slide PowerPoint (presentation.pptx) via python-pptx — no LLM

Public API:
  run(adata, config, study_context, ...) -> ReportWriterResult | None
"""

from __future__ import annotations

import json
import logging
import os
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from ai._base import AiResult
from ai._config_gate import AiDisabledError, check_ai_enabled
from ai._llm_client import call_llm

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Dataclasses
# ---------------------------------------------------------------------------

@dataclass
class ReportSection:
    section_name: str = ""          # "conclusion" | "perspectives"
    section_text: str = ""
    key_findings: list[str] = field(default_factory=list)
    cited_evidence: list[str] = field(default_factory=list)


@dataclass
class ReportWriterResult(AiResult):
    sections: list[ReportSection] = field(default_factory=list)
    pptx_path: str | None = None    # path to written PowerPoint, or None


# ---------------------------------------------------------------------------
# Helpers — parsing
# ---------------------------------------------------------------------------

def _parse_section(raw: str, section_name: str) -> ReportSection:
    """Parse LLM JSON response into a ReportSection. Tolerates partial output."""
    import re

    # Strip markdown fences if present
    cleaned = re.sub(r"```json\s*|```\s*", "", raw).strip()
    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        log.warning("report_writer: JSON parse failed for %s — using raw text", section_name)
        return ReportSection(
            section_name=section_name,
            section_text=raw,
            key_findings=[],
            cited_evidence=[],
        )

    return ReportSection(
        section_name=section_name,
        section_text=data.get("section_text", ""),
        key_findings=data.get("key_findings", []),
        cited_evidence=data.get("cited_evidence", []),
    )


def _section_to_markdown(section: ReportSection) -> str:
    """Render a ReportSection as a markdown string."""
    lines = []
    title = section.section_name.replace("_", " ").title()
    lines.append(f"# {title}\n")
    if section.section_text:
        lines.append(section.section_text)
        lines.append("")
    if section.key_findings:
        lines.append("## Key Findings\n")
        for finding in section.key_findings:
            lines.append(f"- {finding}")
        lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Helpers — PowerPoint
# ---------------------------------------------------------------------------

def _build_pptx(
    study_context: dict,
    analysis_summary: dict,
    sections: list[ReportSection],
    narrative_result: Any | None,
    cluster_annotations: Any | None,
    deg_validation: Any | None,
    output_path: Path,
    figures_dir: Path | None,
) -> None:
    """Build an 8-slide PowerPoint from analysis data. No LLM calls."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx")

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]   # completely blank
    title_layout = prs.slide_layouts[0]   # title + subtitle

    tissue = study_context.get("dataset", {}).get("tissue", "unknown tissue")
    disease = study_context.get("disease", {}).get("context") or "healthy"
    dataset_name = study_context.get("dataset", {}).get("name", "Dataset")
    today = datetime.now(timezone.utc).strftime("%B %Y")

    # Speaker notes helper
    def _add_notes(slide, text: str) -> None:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text

    # Narrative lookup helper
    def _get_narrative_block(block_name: str) -> str:
        if narrative_result is None:
            return ""
        for block in getattr(narrative_result, "blocks", []):
            if getattr(block, "block_name", "") == block_name:
                return getattr(block, "narrative_text", "")
        return ""

    # Figure lookup helper
    def _figure_path(name: str) -> Path | None:
        if figures_dir is None:
            return None
        candidates = [
            figures_dir / f"{name}.png",
            figures_dir / f"{name}.pdf",
            figures_dir / f"{name}.svg",
        ]
        for c in candidates:
            if c.exists():
                return c
        return None

    def _add_title_and_content(slide, title_text: str, content_lines: list[str]) -> None:
        """Add a simple title box and a content text box to a blank slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN  # noqa: F401

        # Title
        txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.runs[0].font.size = Pt(28)
        p.runs[0].font.bold = True

        # Content
        txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        first = True
        for line in content_lines:
            if first:
                tf2.paragraphs[0].text = line
                first = False
            else:
                p2 = tf2.add_paragraph()
                p2.text = line

    # -----------------------------------------------------------------------
    # Slide 1 — Title
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = f"OmicSage Analysis Report\n{dataset_name}"
    slide.placeholders[1].text = f"{tissue.title()} · {disease.title()} · {today}"
    _add_notes(slide, f"Analysis of {dataset_name}: {tissue} tissue, {disease} context. Generated by OmicSage.")

    # -----------------------------------------------------------------------
    # Slide 2 — Data Overview
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    qc = analysis_summary.get("qc_decisions", {})
    sc = analysis_summary.get("study_context", {})
    n_cells = sc.get("n_cells", "N/A")
    n_batches = sc.get("n_batches", "N/A")
    min_genes = qc.get("min_genes", "N/A")
    max_mt = qc.get("max_mt_pct", "N/A")
    removed_pct = qc.get("cells_removed_pct", "N/A")
    modalities = sc.get("modalities", "RNA")
    lines = [
        f"Cells after QC: {n_cells}",
        f"Batches: {n_batches}",
        f"Modalities: {modalities}",
        f"Min genes/cell: {min_genes}",
        f"Max MT%: {max_mt}",
        f"Cells removed by QC: {removed_pct}%",
    ]
    _add_title_and_content(slide, "Data Overview", lines)
    _add_notes(slide, _get_narrative_block("qc_rationale") or "QC metrics for this dataset.")

    # -----------------------------------------------------------------------
    # Slide 3 — UMAP
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    fig = _figure_path("umap") or _figure_path("umap_cell_type")
    if fig is not None:
        try:
            slide.shapes.add_picture(str(fig), Inches(1.5), Inches(1.2), Inches(10), Inches(5.5))
        except Exception:
            _add_title_and_content(slide, "UMAP", ["[Figure could not be loaded]"])
    else:
        _add_title_and_content(slide, "UMAP — Cell Type Embedding", ["[Figure not available]"])
    _add_notes(slide, _get_narrative_block("cell_type_landscape") or "UMAP coloured by cell type.")

    # -----------------------------------------------------------------------
    # Slide 4 — Cell Type Proportions
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    cell_types = analysis_summary.get("cell_types", [])
    if cell_types:
        lines = [f"Cluster {ct.get('cluster', '?')}: {ct.get('cell_type', 'unknown')} "
                 f"({ct.get('n_cells', '?')} cells, {ct.get('confidence', '')})"
                 for ct in cell_types[:12]]
    else:
        lines = ["[Cell type annotations not available]"]
    _add_title_and_content(slide, "Cell Type Proportions", lines)
    _add_notes(slide, _get_narrative_block("cell_type_landscape") or "Cell type proportions.")

    # -----------------------------------------------------------------------
    # Slide 5 — Top DEGs
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    top_degs = analysis_summary.get("top_degs", [])
    if top_degs:
        lines = [f"{d.get('gene', '?')}  log2FC={d.get('log2fc', '?')}  "
                 f"[{d.get('comparison', '')} / cluster {d.get('cluster', '')}]"
                 for d in top_degs[:12]]
    else:
        lines = ["[DEG results not available]"]
    _add_title_and_content(slide, "Top Differentially Expressed Genes", lines)
    _add_notes(slide, _get_narrative_block("differential_expression") or "Top DEGs.")

    # -----------------------------------------------------------------------
    # Slide 6 — Pathway Enrichment
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    pathways = analysis_summary.get("pathways", [])
    if pathways:
        lines = [f"Cluster {p.get('cluster', '?')}: {p.get('pathway', '?')}  "
                 f"(padj={p.get('padj', '?')})"
                 for p in pathways[:12]]
    else:
        lines = ["[Pathway enrichment results not available]"]
    _add_title_and_content(slide, "Pathway Enrichment", lines)
    _add_notes(slide, _get_narrative_block("differential_expression") or "Pathway enrichment.")

    # -----------------------------------------------------------------------
    # Slide 7 — Key Findings
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    conclusion_section = next(
        (s for s in sections if s.section_name == "conclusion"), None
    )
    if conclusion_section and conclusion_section.key_findings:
        findings = conclusion_section.key_findings
    else:
        findings = ["[Key findings not available]"]
    _add_title_and_content(slide, "Key Findings", [f"• {f}" for f in findings])
    _add_notes(slide, conclusion_section.section_text if conclusion_section else "Key findings.")

    # -----------------------------------------------------------------------
    # Slide 8 — Conclusions + Perspectives
    # -----------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    perspectives_section = next(
        (s for s in sections if s.section_name == "perspectives"), None
    )
    conclusions_text = (conclusion_section.section_text[:300] + "…"
                        if conclusion_section and len(conclusion_section.section_text) > 300
                        else (conclusion_section.section_text if conclusion_section else ""))
    perspectives_text = (perspectives_section.section_text[:400] + "…"
                         if perspectives_section and len(perspectives_section.section_text) > 400
                         else (perspectives_section.section_text if perspectives_section else ""))
    lines = []
    if conclusions_text:
        lines += ["Conclusions:", conclusions_text, ""]
    if perspectives_text:
        lines += ["Perspectives:", perspectives_text]
    if not lines:
        lines = ["[Conclusions and perspectives not available]"]
    _add_title_and_content(slide, "Conclusions & Perspectives", lines)
    _add_notes(
        slide,
        (perspectives_section.section_text if perspectives_section else "")
        or _get_narrative_block("interpretation")
        or "Conclusions and perspectives.",
    )

    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def run(
    adata,
    config: dict,
    study_context: dict,
    narrative_result=None,
    coherence_review=None,
    cluster_annotations=None,
    deg_validation=None,
    *,
    report_dir: str,
    figures_dir: str | None = None,
    log_dir: str = "logs/llm",
    runtime_ai: bool = True,
) -> ReportWriterResult | None:
    """
    Generate conclusion + perspectives report sections (LLM) and an
    8-slide PowerPoint (rule-based, no LLM).

    Returns ReportWriterResult or None if AI is disabled.
    """
    try:
        check_ai_enabled(config, module="report_writer", runtime_ai=runtime_ai)
    except AiDisabledError:
        return None

    # -----------------------------------------------------------------------
    # Prepare shared inputs
    # -----------------------------------------------------------------------
    ai_cfg = config.get("ai", {})
    model = ai_cfg.get("model", "llama3")
    provider = ai_cfg.get("provider", "ollama")

    tissue = study_context.get("dataset", {}).get("tissue", "unknown")
    disease_context = study_context.get("disease", {}).get("context") or "not specified"
    biological_question = study_context.get("biological_question", "not specified")
    objectives_raw = study_context.get("objectives", [])
    objectives_str = json.dumps(objectives_raw) if objectives_raw else "not specified"

    # analysis_summary from adata.uns if not passed separately
    analysis_summary_dict: dict = {}
    if adata is not None:
        analysis_summary_dict = getattr(adata, "uns", {}).get("analysis_summary", {})

    # Serialise inputs that are dicts/objects
    analysis_summary_str = json.dumps(analysis_summary_dict) if analysis_summary_dict else "{}"

    ai_narrative_str = "not available"
    if narrative_result is not None:
        blocks = getattr(narrative_result, "blocks", [])
        if blocks:
            ai_narrative_str = "\n\n".join(
                f"## {b.block_name}\n{b.narrative_text}"
                for b in blocks
                if getattr(b, "narrative_text", "")
            )

    coherence_str = "not available"
    if coherence_review is not None:
        try:
            import dataclasses
            coherence_str = json.dumps(dataclasses.asdict(coherence_review))
        except Exception:
            coherence_str = str(coherence_review)

    shared_inputs = dict(
        tissue=tissue,
        disease_context=disease_context,
        biological_question=biological_question,
        objectives=objectives_str,
        analysis_summary=analysis_summary_str,
        ai_narrative=ai_narrative_str,
        coherence_review=coherence_str,
    )

    # -----------------------------------------------------------------------
    # Two LLM calls — one per section
    # -----------------------------------------------------------------------
    sections: list[ReportSection] = []
    last_raw = ""
    last_reasoning = ""

    for section_name in ("conclusion", "perspectives"):
        inputs = {**shared_inputs, "report_section": section_name}
        try:
            raw = call_llm(
                skill_name="report_writer",
                inputs=inputs,
                config=config,
                log_dir=log_dir,
                module="report_writer",
                runtime_ai=runtime_ai,
            )
        except Exception as exc:
            log.warning("report_writer: LLM call failed for %s — %s", section_name, exc)
            raw = "{}"

        section = _parse_section(raw, section_name)
        sections.append(section)
        last_raw = raw
        last_reasoning = section.section_name  # just a marker; real reasoning in section

    # Extract reasoning from the last successful parse
    reasoning = ""
    try:
        import re
        cleaned = re.sub(r"```json\s*|```\s*", "", last_raw).strip()
        reasoning = json.loads(cleaned).get("reasoning", "")
    except Exception:
        pass

    # -----------------------------------------------------------------------
    # Write markdown files
    # -----------------------------------------------------------------------
    report_path = Path(report_dir)
    report_path.mkdir(parents=True, exist_ok=True)

    for section in sections:
        md_path = report_path / f"{section.section_name}.md"
        md_path.write_text(_section_to_markdown(section), encoding="utf-8")
        log.info("report_writer: wrote %s", md_path)

    # -----------------------------------------------------------------------
    # Build PowerPoint (soft dependency)
    # -----------------------------------------------------------------------
    pptx_path: str | None = None
    figures_path = Path(figures_dir) if figures_dir else None

    try:
        pptx_out = report_path / "presentation.pptx"
        _build_pptx(
            study_context=study_context,
            analysis_summary=analysis_summary_dict,
            sections=sections,
            narrative_result=narrative_result,
            cluster_annotations=cluster_annotations,
            deg_validation=deg_validation,
            output_path=pptx_out,
            figures_dir=figures_path,
        )
        pptx_path = str(pptx_out)
        log.info("report_writer: wrote %s", pptx_out)
    except ImportError:
        log.warning("report_writer: python-pptx not installed — skipping PowerPoint")
    except Exception as exc:
        log.warning("report_writer: PowerPoint generation failed — %s", exc)

    # -----------------------------------------------------------------------
    # Assemble result
    # -----------------------------------------------------------------------
    return ReportWriterResult(
        timestamp=datetime.now(timezone.utc).isoformat(),
        model=model,
        provider=provider,
        skill_name="report_writer",
        skill_version="1.0",
        reasoning=reasoning,
        sections=sections,
        pptx_path=pptx_path,
    )
